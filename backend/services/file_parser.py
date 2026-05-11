"""
Universal File Parser
Supports: PDF, DOCX, PPTX, XLSX, CSV, TXT, MD, JSON, XML, HTML, Images (OCR)
"""

import io
import json
import csv
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import Optional
from loguru import logger


# ── Supported file types ────────────────────────────────────
SUPPORTED_EXTENSIONS = {
    ".pdf", ".docx", ".doc",
    ".pptx", ".ppt",
    ".xlsx", ".xls", ".csv",
    ".txt", ".md", ".markdown",
    ".json", ".xml", ".html", ".htm",
    ".png", ".jpg", ".jpeg", ".tiff", ".bmp", ".webp",
}


def get_file_type(filename: str) -> str:
    return Path(filename).suffix.lower()


def is_supported(filename: str) -> bool:
    return get_file_type(filename) in SUPPORTED_EXTENSIONS


# ── Main Parser ─────────────────────────────────────────────
def parse_file(file_bytes: bytes, filename: str) -> dict:
    """
    Parse any supported file and return:
    {
        "text": str,           # extracted full text
        "metadata": dict,      # page counts, sheet names, etc.
        "file_type": str,      # e.g. ".pdf"
        "error": str | None    # error message if failed
    }
    """
    ext = get_file_type(filename)
    logger.info(f"Parsing file: {filename} (type: {ext})")

    parsers = {
        ".pdf":      _parse_pdf,
        ".docx":     _parse_docx,
        ".doc":      _parse_docx,
        ".pptx":     _parse_pptx,
        ".ppt":      _parse_pptx,
        ".xlsx":     _parse_excel,
        ".xls":      _parse_excel,
        ".csv":      _parse_csv,
        ".txt":      _parse_text,
        ".md":       _parse_text,
        ".markdown": _parse_text,
        ".json":     _parse_json,
        ".xml":      _parse_xml,
        ".html":     _parse_html,
        ".htm":      _parse_html,
        ".png":      _parse_image,
        ".jpg":      _parse_image,
        ".jpeg":     _parse_image,
        ".tiff":     _parse_image,
        ".bmp":      _parse_image,
        ".webp":     _parse_image,
    }

    parser = parsers.get(ext)
    if not parser:
        return {
            "text": "",
            "metadata": {},
            "file_type": ext,
            "error": f"Unsupported file type: {ext}"
        }

    try:
        result = parser(file_bytes, filename)
        result["file_type"] = ext
        result["error"] = None
        logger.info(f"Successfully parsed {filename}: {len(result.get('text',''))} chars")
        return result
    except Exception as e:
        logger.error(f"Error parsing {filename}: {e}")
        return {
            "text": "",
            "metadata": {},
            "file_type": ext,
            "error": str(e)
        }


# ── PDF ─────────────────────────────────────────────────────
def _parse_pdf(file_bytes: bytes, filename: str) -> dict:
    import fitz  # PyMuPDF
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    pages_text = []
    for page_num, page in enumerate(doc, start=1):
        text = page.get_text("text").strip()
        if text:
            pages_text.append(f"[Page {page_num}]\n{text}")
    return {
        "text": "\n\n".join(pages_text),
        "metadata": {
            "page_count": len(doc),
            "title": doc.metadata.get("title", ""),
            "author": doc.metadata.get("author", ""),
        }
    }


# ── DOCX ────────────────────────────────────────────────────
def _parse_docx(file_bytes: bytes, filename: str) -> dict:
    from docx import Document
    doc = Document(io.BytesIO(file_bytes))
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]

    # Extract tables too
    tables_text = []
    for i, table in enumerate(doc.tables):
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(" | ".join(cells))
        tables_text.append(f"[Table {i+1}]\n" + "\n".join(rows))

    full_text = "\n\n".join(paragraphs)
    if tables_text:
        full_text += "\n\n" + "\n\n".join(tables_text)

    return {
        "text": full_text,
        "metadata": {
            "paragraph_count": len(paragraphs),
            "table_count": len(doc.tables),
        }
    }


# ── PPTX ────────────────────────────────────────────────────
def _parse_pptx(file_bytes: bytes, filename: str) -> dict:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(file_bytes))
    slides_text = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_content = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_content.append(shape.text.strip())
        if slide_content:
            slides_text.append(f"[Slide {i}]\n" + "\n".join(slide_content))
    return {
        "text": "\n\n".join(slides_text),
        "metadata": {"slide_count": len(prs.slides)}
    }


# ── XLSX / XLS ──────────────────────────────────────────────
def _parse_excel(file_bytes: bytes, filename: str) -> dict:
    import pandas as pd
    xls = pd.ExcelFile(io.BytesIO(file_bytes))
    sheets_text = []
    for sheet_name in xls.sheet_names:
        df = xls.parse(sheet_name).fillna("")
        sheet_str = f"[Sheet: {sheet_name}]\n{df.to_string(index=False)}"
        sheets_text.append(sheet_str)
    return {
        "text": "\n\n".join(sheets_text),
        "metadata": {"sheet_names": xls.sheet_names, "sheet_count": len(xls.sheet_names)}
    }


# ── CSV ─────────────────────────────────────────────────────
def _parse_csv(file_bytes: bytes, filename: str) -> dict:
    import pandas as pd
    df = pd.read_csv(io.BytesIO(file_bytes)).fillna("")
    return {
        "text": df.to_string(index=False),
        "metadata": {
            "row_count": len(df),
            "column_count": len(df.columns),
            "columns": list(df.columns),
        }
    }


# ── TXT / MD ────────────────────────────────────────────────
def _parse_text(file_bytes: bytes, filename: str) -> dict:
    text = file_bytes.decode("utf-8", errors="ignore")
    return {
        "text": text,
        "metadata": {
            "char_count": len(text),
            "line_count": text.count("\n"),
        }
    }


# ── JSON ────────────────────────────────────────────────────
def _parse_json(file_bytes: bytes, filename: str) -> dict:
    data = json.loads(file_bytes.decode("utf-8", errors="ignore"))
    # Pretty print JSON as readable text
    text = json.dumps(data, indent=2, ensure_ascii=False)
    return {
        "text": text,
        "metadata": {"type": type(data).__name__}
    }


# ── XML ─────────────────────────────────────────────────────
def _parse_xml(file_bytes: bytes, filename: str) -> dict:
    root = ET.fromstring(file_bytes.decode("utf-8", errors="ignore"))
    texts = [elem.text.strip() for elem in root.iter() if elem.text and elem.text.strip()]
    return {
        "text": "\n".join(texts),
        "metadata": {"root_tag": root.tag}
    }


# ── HTML ────────────────────────────────────────────────────
def _parse_html(file_bytes: bytes, filename: str) -> dict:
    from bs4 import BeautifulSoup
    soup = BeautifulSoup(file_bytes, "html.parser")
    # Remove script and style tags
    for tag in soup(["script", "style", "nav", "footer"]):
        tag.decompose()
    text = soup.get_text(separator="\n").strip()
    # Clean up excessive blank lines
    lines = [l.strip() for l in text.splitlines() if l.strip()]
    return {
        "text": "\n".join(lines),
        "metadata": {"title": soup.title.string if soup.title else ""}
    }


# ── Images (OCR) ────────────────────────────────────────────
def _parse_image(file_bytes: bytes, filename: str) -> dict:
    try:
        import pytesseract
        from PIL import Image
        image = Image.open(io.BytesIO(file_bytes))
        text = pytesseract.image_to_string(image)
        return {
            "text": text.strip(),
            "metadata": {
                "width": image.width,
                "height": image.height,
                "mode": image.mode,
                "ocr": True,
            }
        }
    except Exception as e:
        # Tesseract not installed — degrade gracefully
        logger.warning(f"OCR failed for {filename}: {e}")
        return {
            "text": f"[Image file: {filename} — OCR unavailable. Install Tesseract to extract text.]",
            "metadata": {"ocr": False, "error": str(e)}
        }
